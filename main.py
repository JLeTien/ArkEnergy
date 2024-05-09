from scrapy import signals
from scrapy.crawler import CrawlerProcess
from scrapy.utils.project import get_project_settings
from ArkEnergyScraper.Scraper.spiders.ArkSpiders import Spider1
from ArkEnergyScraper.Scraper.spiders.ArkSpiders import Spider2
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import os
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR 
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from pptx.dml.color import RGBColor
import yfinance as yf
import pandas as pd

def slide1(prs):
    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[5]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "ARK Energy"
    title.top = Inches(2.5)
    title.left = Inches(1)
    title.width = Inches(2.66)
    title.height = Inches(1.42)
    
    current_time= slide.shapes.add_textbox(left=Inches(0.9055), top=Inches(4.114), width=Inches(2.838), height=Inches(0.4055))
    current_datetime = datetime.now().strftime("%Y-%m-%d")
    # to generate time stamps use formating "%H:%M:%S" 
    current_time.text = f"Generated at: {current_datetime}"

    # Get the current directory of the Python script
    current_dir = os.path.dirname(os.path.abspath(__file__))
    img_folder = "Logo"

    # Path to the image file relative to the current directory
    img1 = "180degreelogo.png"
    img_path1 = os.path.join(current_dir, img_folder, img1)
    from_left = Inches(7.456)
    from_top = Inches(6.5)
    length = Inches(2)
    slide.shapes.add_picture(img_path1, from_left, from_top, length)

    img2 = "Arkenergylogo.png"
    img_path2 = os.path.join(current_dir, img_folder, img2)
    from_left = Inches(5.169)
    from_top = Inches(6.5)
    length = Inches(2)
    slide.shapes.add_picture(img_path2, from_left, from_top, length)
    
    img3 = "Home_button.png"
    img_path3 = os.path.join(current_dir, img_folder, img3)
    from_left = Inches(8.413)
    from_top = Inches(0.0511)
    height = Inches(0.5)
    width = Inches(0.5)
    add_picture3 = slide.shapes.add_picture(img_path3, from_left, from_top, height, width)
    add_picture3.zorder = 100
    
    green_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.921), Inches(1), Inches(3.921), Inches(6.212))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    
    # Home button 

    home_button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(8.311),top=Inches(0),width=Inches(0.708),height=Inches(0.606))
    fill_home = home_button.fill
    fill_home.solid()
    fill_home.fore_color.rgb = RGBColor(148, 197, 84)
    fill_home.zorder = 1
    line_home = home_button.line
    line_home.color.rgb = RGBColor(148, 197, 84)  
    line_home.zorder = -16

    # COMMODITY BUTTON

    box_commodities = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_commodities = box_commodities.fill
    fill_commodities.solid()
    fill_commodities.fore_color.rgb = RGBColor(148, 197, 84)
    fill_commodities.zorder = 1
    line = box_commodities.line
    line.color.rgb = RGBColor(148, 197, 84)  # HEX: 94C554

    textbox_c = slide.shapes.add_textbox(left=Inches(4.295), top=Inches(2.307), height=Inches(0.41), width=Inches(1.58))
    text_box_c= textbox_c.text_frame
    text_box_c.text = "Commodities"

    # NEWS BUTTON 

    box_news = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_news = box_news.fill
    fill_news.solid()
    fill_news.fore_color.rgb = RGBColor(148, 197, 84)
    fill_news.zorder = 1
    line_news = box_news.line
    line_news.color.rgb = RGBColor(148, 197, 84)  

    textbox_n = slide.shapes.add_textbox(left=Inches(6.464), top=Inches(2.311), height=Inches(0.41), width=Inches(0.8))
    text_box_n= textbox_n.text_frame
    text_box_n.text = "News"

    # COMPETITORS/OMES BUTTON

    box_comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(7.877),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_comp = box_comp.fill
    fill_comp.solid()
    fill_comp.fore_color.rgb = RGBColor(148, 197, 84)
    fill_comp.zorder = 2
    line_comp = box_comp.line
    line_comp.color.rgb = RGBColor(148, 197, 84) 

    textbox_comp = slide.shapes.add_textbox(left=Inches(7.944), top=Inches(2.314), height=Inches(0.41), width=Inches(1.527))
    text_box_comp= textbox_comp.text_frame
    text_box_comp.text = "Competitors"

    # PROJECTS BUTTON

    box_projects = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_projects = box_projects.fill
    fill_projects.solid()
    fill_projects.fore_color.rgb = RGBColor(148, 197, 84)
    fill_projects.zorder = 2
    line_projects = box_projects.line
    line_projects.color.rgb = RGBColor(148, 197, 84) 

    textbox_p = slide.shapes.add_textbox(left=Inches(4.551), top=Inches(3.543), height=Inches(0.41), width=Inches(1.027))
    text_box_p= textbox_p.text_frame
    text_box_p.text = "Projects"

    # GRANTS BUTTON

    box_grants = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 2
    line_grants = box_grants.line
    line_grants.color.rgb = RGBColor(148, 197, 84) 

    textbox_g = slide.shapes.add_textbox(left=Inches(6.417), top=Inches(3.543), height=Inches(0.41), width=Inches(0.901))
    text_box_g= textbox_g.text_frame
    text_box_g.text = "Grants"

def slide2(prs, slide_data):
    # Add a slide to the presentation
    slide_layout = prs.slide_layouts[1]  # Choose a layout that supports tables (e.g., Title and Content)
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title of the slide
    title_shape = slide.shapes.title
    title_shape.text = "Recent News"
    
    # Define table dimensions and position
    rows, cols = 2, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1.5)

    # Add a table to the slide
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # Populate the table with data
    for col, slide_content in enumerate(slide_data):
        table.cell(0, col).text = slide_content['title']
        table.cell(1, col).text = slide_content['content']

    # Set the font size for the table cells
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(12)  

def slide3(prs, slide_data):
    # Add a slide to the presentation
    slide_layout = prs.slide_layouts[1]  # Choose a layout that supports tables (e.g., Title and Content)
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title of the slide
    title_shape = slide.shapes.title
    title_shape.text = "Recent News"
    
    # Define table dimensions and position
    rows, cols = 2, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1.5)

    # Add a table to the slide
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # Populate the table with data
    for col, slide_content in enumerate(slide_data):
        table.cell(0, col).text = slide_content['title']
        table.cell(1, col).text = slide_content['content']

    # Set the font size for the table cells
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(12)  

def slide4(prs, company_data):
    # Add a slide to the presentation
    slide_layout = prs.slide_layouts[5]  # Choose a layout that supports title and content
    slide = prs.slides.add_slide(slide_layout)

    # Set the title of the slide
    title_shape = slide.shapes.title
    title_shape.text = "Company Data"
    title_shape.text_frame.paragraphs[0].font.size = Pt(15)

    # Define table position and size
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    left = Inches(0.1)  # Adjust left position to center the table
    top = Inches(1)
    width = slide_width - Inches(1)  # Use slide width minus some padding
    height = slide_height - Inches(3)  # Use slide height minus some padding

    # Add a table to the slide
    table = slide.shapes.add_table(len(company_data) + 1, len(company_data[0]), left, top, width, height).table

    # Populate the table with data
    headers = list(company_data[0].keys())
    for col, header in enumerate(headers):
        table.cell(0, col).text = header

    for row_idx, row_data in enumerate(company_data, start=1):
        for col_idx, value in enumerate(row_data.values()):
            table.cell(row_idx, col_idx).text = str(value)

    # Set the font size for the table cells
    for row in range(len(company_data) + 1):
        for col in range(len(company_data[0])):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

def slide5(prs):
    slide_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(slide_layout)
    title4 = slide4.shapes.title
    title4.text = "Hyperlinks"
    para4 = slide4.placeholders[1].text_frame.paragraphs[0]
    addrun1 = para4.add_run()
    addrun1.text = "Google Hyperlink"
    hlink1 = addrun1.hyperlink
    hlink1.address = "https://www.google.com.au"
    
def generate_ppt(slides_data, slides_data2, company_data):
    prs = Presentation()
    slide1(prs)    
    slide2(prs, slides_data)
    slide3(prs, slides_data2)
    slide4(prs, company_data)
    slide5(prs)
    prs.save('Monthly_Report.pptx')

def main():
    process = CrawlerProcess()
    process.crawl(Spider1)
    process.crawl(Spider2)
    process.start()

    slides_data = Spider1.slides_data
    slides_data2 = Spider2.slides_data
    
    # Now you can use the slides_data list as needed
    # print(slides_data)
    company_data = pd.read_csv("company_data.csv").to_dict("records")
    
    generate_ppt(slides_data, slides_data2, company_data)
    
if __name__ == "__main__":
    main()