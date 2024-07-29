from scrapy.crawler import CrawlerProcess
from ArkEnergyScraper.Scraper.spiders.ArkSpiders import Spider1
from ArkEnergyScraper.Scraper.spiders.ArkSpiders import Spider2
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import os
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from datetime import datetime, timedelta
from pptx.dml.color import RGBColor
import yfinance as yf
import pandas as pd
from pptx.enum.text import MSO_ANCHOR
from AI_module.summarygemini import *

def title_slide(prs):
    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "ARK Energy"
    title.top = Inches(2.5)
    title.left = Inches(1)
    title.width = Inches(2.66)
    title.height = Inches(1.42)
    
    current_time= slide.shapes.add_textbox(left=Inches(0.9055), top=Inches(4.114), width=Inches(2.838), height=Inches(0.4055))
    current_datetime = datetime.now().strftime("%Y-%m-%d")
    current_time.text = f"Generated at: {current_datetime}"

    # Get the current directory of the Python script
    current_dir = os.path.dirname(os.path.abspath(__file__))
    img_folder = "Logo"

    # Path to the image file relative to the current directory
    img1 = "180degreelogo.png"
    img_path1 = os.path.join(current_dir, img_folder, img1)
    from_left = Inches(7.456)
    from_top = Inches(6.5)
    length = Inches(2)
    slide.shapes.add_picture(img_path1, from_left, from_top, length)

    img2 = "Arkenergylogo.png"
    img_path2 = os.path.join(current_dir, img_folder, img2)
    from_left = Inches(5.169)
    from_top = Inches(6.5)
    length = Inches(2)
    slide.shapes.add_picture(img_path2, from_left, from_top, length)
    
    img3 = "Home_button.png"
    img_path3 = os.path.join(current_dir, img_folder, img3)
    from_left = Inches(8.413)
    from_top = Inches(0.0511)
    height = Inches(0.5)
    width = Inches(0.5)
    add_picture3 = slide.shapes.add_picture(img_path3, from_left, from_top, height, width)
    add_picture3.zorder = 100
    
    green_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.921), Inches(1), Inches(3.921), Inches(6.212))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    
    # Home button 
    home_button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(8.311),top=Inches(0),width=Inches(0.708),height=Inches(0.606))
    fill_home = home_button.fill
    fill_home.solid()
    fill_home.fore_color.rgb = RGBColor(148, 197, 84)
    fill_home.zorder = 1
    line_home = home_button.line
    line_home.color.rgb = RGBColor(148, 197, 84)  
    line_home.zorder = -16

    # COMMODITY BUTTON
    box_commodities = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_commodities = box_commodities.fill
    fill_commodities.solid()
    fill_commodities.fore_color.rgb = RGBColor(148, 197, 84)
    fill_commodities.zorder = 1
    line = box_commodities.line
    line.color.rgb = RGBColor(148, 197, 84)  # HEX: 94C554

    textbox_c = slide.shapes.add_textbox(left=Inches(4.295), top=Inches(2.307), height=Inches(0.41), width=Inches(1.58))
    text_box_c= textbox_c.text_frame
    text_box_c.text = "Commodities"

    # NEWS BUTTON 
    box_news = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_news = box_news.fill
    fill_news.solid()
    fill_news.fore_color.rgb = RGBColor(148, 197, 84)
    fill_news.zorder = 1
    line_news = box_news.line
    line_news.color.rgb = RGBColor(148, 197, 84)  

    textbox_n = slide.shapes.add_textbox(left=Inches(6.464), top=Inches(2.311), height=Inches(0.41), width=Inches(0.8))
    text_box_n= textbox_n.text_frame
    text_box_n.text = "News"

    # COMPETITORS/OMES BUTTON
    box_comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(7.877),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_comp = box_comp.fill
    fill_comp.solid()
    fill_comp.fore_color.rgb = RGBColor(148, 197, 84)
    fill_comp.zorder = 2
    line_comp = box_comp.line
    line_comp.color.rgb = RGBColor(148, 197, 84) 

    textbox_comp = slide.shapes.add_textbox(left=Inches(7.944), top=Inches(2.314), height=Inches(0.41), width=Inches(1.527))
    text_box_comp= textbox_comp.text_frame
    text_box_comp.text = "Competitors"

    # PROJECTS BUTTON
    box_projects = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_projects = box_projects.fill
    fill_projects.solid()
    fill_projects.fore_color.rgb = RGBColor(148, 197, 84)
    fill_projects.zorder = 2
    line_projects = box_projects.line
    line_projects.color.rgb = RGBColor(148, 197, 84) 

    textbox_p = slide.shapes.add_textbox(left=Inches(4.551), top=Inches(3.543), height=Inches(0.41), width=Inches(1.027))
    text_box_p= textbox_p.text_frame
    text_box_p.text = "Projects"

    # GRANTS BUTTON
    box_grants = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 2
    line_grants = box_grants.line
    line_grants.color.rgb = RGBColor(148, 197, 84) 

    textbox_g = slide.shapes.add_textbox(left=Inches(6.417), top=Inches(3.543), height=Inches(0.41), width=Inches(0.901))
    text_box_g= textbox_g.text_frame
    text_box_g.text = "Grants"

def slide000(prs, company_data):
    # Add a slide to the presentation
    slide_layout = prs.slide_layouts[5]  # Choose a layout that supports title and content
    slide = prs.slides.add_slide(slide_layout)

    # Set the title of the slide
    title_shape = slide.shapes.title
    title_shape.text = "Company Data"
    title_shape.text_frame.paragraphs[0].font.size = Pt(15)

    # Define table position and size
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    left = Inches(0.1)  # Adjust left position to center the table
    top = Inches(1)
    width = slide_width - Inches(1)  # Use slide width minus some padding
    height = slide_height - Inches(3)  # Use slide height minus some padding

    # Add a table to the slide
    table = slide.shapes.add_table(len(company_data) + 1, len(company_data[0]), left, top, width, height).table

    # Populate the table with data
    headers = list(company_data[0].keys())
    for col, header in enumerate(headers):
        table.cell(0, col).text = header

    for row_idx, row_data in enumerate(company_data, start=1):
        for col_idx, value in enumerate(row_data.values()):
            table.cell(row_idx, col_idx).text = str(value)

    # Set the font size for the table cells
    for row in range(len(company_data) + 1):
        for col in range(len(company_data[0])):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

def commodities_slide(prs):
    # Slide 2: Content slide with bullet points
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "Key Commodities"
    title2.top = Inches(0.1)
    title2.left = Inches(0.3)
    title2.width = Inches(3)
    title2.height = Inches(1.42)
    title2.text_frame.paragraphs[0].font.size = Pt(26)

    # ONE TABLE 
    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(6)
    second = slide2.shapes.add_table(8, 4, x, y, cx, cy)
    second_table = second.table
    for row in second_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 
    cell = second_table.cell(0, 1)
    cell.text = 'Gold'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell = second_table.cell(0, 2)
    cell.text = 'Oil'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell = second_table.cell(0, 3)
    cell.text = 'Crude Oil'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    current_datetime = "2024-04-30"
    labels = [
         "Chart (1 year)",
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares"
    ]
    for i, label in enumerate(labels):
        cell = second_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
def news_slide(prs, slides_data):
    # Slide 3: Content slide with image
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide3 = prs.slides.add_slide(slide_layout)
    title3 = slide3.shapes.title
    title3.text = "News"

    # Defining the number of rows and columns for the table
    rows = 4
    cols = 4

    # Define the height of each row
    row_heights = [Inches(0.8), Inches(0.6), Inches(0.6), Inches(3)]

    # Adding a table 
    left = Inches(0.5)
    top = Inches(1.3)
    total_width = Inches(9)
    height = Inches(5.5)
    
    # Create the table shape
    table = slide3.shapes.add_table(rows, cols, left, top, total_width, height).table

    # Set the width for each column
    col_widths = [Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.5)]
    
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Define colors
    green_color = RGBColor(148, 197, 84)
    white_color = RGBColor(255, 255, 255)

    # Data to populate the table
    data = [
        ['', slides_data[0]['title'], slides_data[1]['title'], slides_data[2]['title']],
        ['Date', 'xx/xx/xxxx', 'xx/xx/xxxx', 'xx/xx/xxxx'],
        ['Source', 'URL', 'URL', 'URL'],
        ['Description', summarise_gemini(slides_data[0]['content']), summarise_gemini(slides_data[1]['content']), summarise_gemini(slides_data[2]['content'])]
    ]

    # Populate the table with data and formatting
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
            
            # Set background color
            cell.fill.solid()
            if i == 0 or j == 0:
                cell.fill.fore_color.rgb = green_color
                cell.text_frame.paragraphs[0].font.size = Pt(14)  # Set font size for headers
            else:
                cell.fill.fore_color.rgb = white_color
                if i == 3:  # Description row
                    cell.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size for description content
                else:
                    cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size for other content

    # Adjust row heights
    for row_idx, height in enumerate(row_heights):
        table.rows[row_idx].height = height

def comp_slide(prs, financial_data):
    slide_layout = prs.slide_layouts[5]
    slide4 = prs.slides.add_slide(slide_layout)
    title4 = slide4.shapes.title
    title4.text = "Competitors/OEMs"
    title4.top = Inches(0.05)
    title4.left = Inches(0.35)
    title4.width = Inches(3)
    title4.height = Inches(1.42)
    title4.text_frame.paragraphs[0].font.size = Pt(25)
    
    # Define the number of rows and columns
    num_rows = 6  # Number of metrics + header row
    num_cols = len(financial_data) + 1  # Number of companies + 1 for the row labels

    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9.2), Inches(6)
    table = slide4.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table

    # Define colors
    header_color = RGBColor(148, 197, 84)
    white_color = RGBColor(255, 255, 255)

    # Define headers
    metrics = [
        "Company Name",
        "Current Price",
        "Price 1 Week Ago",
        "Price 1 Month Ago",
        "Change in Price (%) Week on Week",
        "Change in Price (%) 1 Month Ago"
    ]

    # Populate the header row
    for col_index, metric in enumerate(metrics):
        cell = table.cell(0, col_index)
        cell.text = metric
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color  # Header color

    # Populate the company names in the first column
    for row_index, data in enumerate(financial_data, start=1):
        cell = table.cell(row_index, 0)
        cell.text = data["Company Name"]
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color  # First column color

    # Populate the data in the table and set cell colors
    for row_index, metric in enumerate(metrics[1:], start=1):
        for col_index, data in enumerate(financial_data, start=1):
            value = data.get(
                {
                    "Current Price": "Current Price",
                    "Price 1 Week Ago": "Price 1 Week Ago",
                    "Price 1 Month Ago": "Price 1 Month Ago",
                    "Change in Price (%) Week on Week": "Change in Price (%) Week on Week",
                    "Change in Price (%) 1 Month Ago": "Change in Price (%) 1 Month Ago"
                }[metric], "N/A"
            )
            if value is None:
                value = "N/A"
            elif isinstance(value, (int, float)):
                value = f"{value:.2f}"
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = white_color
   
def project_slide(prs):
    # Slide 5: Projects
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide5 = prs.slides.add_slide(slide_layout)
    title5 = slide5.shapes.title
    title5.text = "Projects"
    title5.top = Inches(0.08)
    title5.left = Inches(0.45)
    title5.width = Inches(1.72)
    title5.height = Inches(1.42)
    title5.text_frame.paragraphs[0].font.size = Pt(28)

    rows = 6
    cols = 2
    
    # Define the width of each column
    col_widths = [Inches(2)] * cols
    
    # Define the height of each row
    row_heights = [Inches(1), Inches(1), Inches(1), Inches(1)]
   
    # Adding a table 
    left = Inches(0.45)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(12)
    table = slide5.shapes.add_table(rows, cols, left, top, width, height).table
   
    # Set alternating fill colors for the rows
    for row_idx, row in enumerate(table.rows):
        if row_idx % 2 == 0:
            color = RGBColor(148, 197, 84)  # Green
            row.height = row_heights[0]
        else:
            color = RGBColor(255, 255, 255)  # White
            row.height = row_heights[1]
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = color
    
    # Add data to the table
    data = [
        ['A1', 'B1'],
        ['A2', 'B2'],
        ['A3', 'B3'],
        ['A4', 'B4']
    ]
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
            cell.text_frame.paragraphs[0].font.size = Pt(12)  
    
    # Project Heading 
    cell = table.cell(0, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(1, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(2, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(3, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(4, 0)
    cell.text = 'XX/XX/XXXX Title'
    
    # Content of Projects 
    cell = table.cell(0, 1)
    cell.text = 'Content'
    cell = table.cell(1, 1)
    cell.text = 'Content'
    cell = table.cell(2, 1)
    cell.text = 'Content'
    cell = table.cell(3, 1)
    cell.text = 'Content'
    cell = table.cell(4, 1)
    cell.text = 'Content'
    first = table.columns[0]
    first.width = Inches(2.5)
    second = table.columns[1]
    second.width = Inches(6.5)

def grants_slide(prs):
    # Slide 6: Adding hyperlinks
    slide_layout = prs.slide_layouts[5] 
    slide6 = prs.slides.add_slide(slide_layout)
    title6 = slide6.shapes.title
    title6.text = "Grants"
    title6.top = Inches(0.08)
    title6.left = Inches(0.45)
    title6.width = Inches(1.72)
    title6.height = Inches(1.42)
    title6.text_frame.paragraphs[0].font.size = Pt(28)
   
    # Grant 1
    # Big White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(1.24),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)
    # Green Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(3.40),width=Inches(4.055),height=Inches(0.54))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(3.67),top=Inches(3.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text 
    textbox= slide6.shapes.add_textbox(left=Inches(0.84), top=Inches(1.4), height=Inches(0.41), width=Inches(1.58))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(3.67), top=Inches(3.468), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(1.22), Inches(0.7), Inches(3.94))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
    
    # GRANT 2 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.248),top=Inches(1.24),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.248),top=Inches(3.40),width=Inches(4.055),height=Inches(0.54))
    fill_apply = box_apply.fill
    fill_apply.solid()
    fill_apply.fore_color.rgb = RGBColor(148, 197, 84)
    fill_apply.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(8.2),top=Inches(3.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text
    textbox= slide6.shapes.add_textbox(left=Inches(5.38), top=Inches(1.37), height=Inches(0.41), width=Inches(1.58))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(8.2), top=Inches(3.46), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.24), Inches(1.22), Inches(5.24), Inches(3.94))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
   
    # GRANT 3 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(4.22),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(6.39),width=Inches(4.055),height=Inches(0.54))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(3.67),top=Inches(6.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text
    textbox= slide6.shapes.add_textbox(left=Inches(0.84), top=Inches(4.53), height=Inches(1.31), width=Inches(3.91))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(3.67), top=Inches(6.448), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(4.22), Inches(0.7), Inches(6.93))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
    
    # GRANT 4 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.24),top=Inches(4.22),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.24),top=Inches(6.39),width=Inches(4.055),height=Inches(0.54))
    fill_apply = box_apply.fill
    fill_apply.solid()
    fill_apply.fore_color.rgb = RGBColor(148, 197, 84)
    fill_apply.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(8.2),top=Inches(6.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Text Box
    textbox= slide6.shapes.add_textbox(left=Inches(5.38), top=Inches(4.53), height=Inches(1.31), width=Inches(3.91))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(8.2), top=Inches(6.448), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.24), Inches(4.2), Inches(5.24), Inches(6.93))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
    
def disclaimer_slide(prs):
    # Slide 7 Disclaimer
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide7 = prs.slides.add_slide(slide_layout)
    title7 = slide7.shapes.title
    title7.text = "Disclaimer"
    numbered_list_items = [
        "Information Collection: The data collected through this project is sourced from publicly available websites and sources related to the ammonia and renewable energy sectors in Australia. It is intended to provide insights and analysis into these industries.",
        "Non-Commercial Use: The data and insights obtained from this project are not to be used for commercial purposes, including but not limited to selling or distributing the data for profit without explicit authorization.",
        "Accuracy and Reliability: While efforts are made to ensure the accuracy and reliability of the data collected, we do not guarantee its completeness or correctness. Users should independently verify the data before making any decisions or relying on it for business or investment purposes.",
        "No Endorsement: The inclusion of any specific data or information in this project does not constitute an endorsement or recommendation of any particular company, product, or service mentioned in the scraped data.",
        "No Liability: We disclaim any liability for damages, losses, or legal consequences arising from the use or misuse of the data collected through this project. Users assume all risks associated with their use of the information."
    ]
    # Add a text box for the numbered list
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6.5)
    text_box = slide7.shapes.add_textbox(left, top, width, height)
    # Set text alignment
    text_frame = text_box.text_frame
    text_frame.text_anchor = MSO_ANCHOR.TOP
    # Allow text to wrap within the text box
    text_frame.word_wrap = True
    # Add paragraphs with numbered text
    for idx, item in enumerate(numbered_list_items, start=1):
        p = text_frame.add_paragraph()
        p.text = f"{idx}. {item}"
        p.alignment = PP_ALIGN.LEFT
        # Set the paragraph to be numbered
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.level = 0  # Set the level of the numbered list (0 is the first level)
    # Add the disclaimer paragraph
    disclaimer = ("By using this program, you acknowledge that you have read, understood, "
                "and agree to be bound by this disclaimer.")
    p = text_frame.add_paragraph()
    p.text = disclaimer
    p.alignment = PP_ALIGN.LEFT
    
def references_slide(prs):
     # Slide 8 Adding hyperlinks
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide8 = prs.slides.add_slide(slide_layout)
    title8 = slide8.shapes.title
    title8.text = "References"
    para8 = slide8.placeholders[1].text_frame.paragraphs[0]
    addrun1 = para8.add_run()
    addrun1.text = "Google Hyperlink"
    hlink1 = addrun1.hyperlink
    hlink1.address = "https://www.google.com.au"

def get_financial_data():
    symbols = {
        "Hyzon Motors Inc.": "HYZN",
        "Energy Vault Holdings, Inc.": "NRGV",
        "Hyundai Motor Company": "HYMTF",
        "Toyota Motor Corporation": "TM",
        "Cummins Inc.": "CMI"
    }

    end_date = datetime.now()
    start_date_week = end_date - timedelta(weeks=1)
    start_date_month = end_date - timedelta(weeks=4)

    data = []

    for company_name, symbol in symbols.items():
        stock = yf.Ticker(symbol)
        hist = stock.history(period="1mo")
        
        if not hist.empty:
            current_price = hist['Close'].iloc[-1]
            price_week_ago = hist.loc[hist.index >= start_date_week.strftime('%Y-%m-%d')]['Close'].iloc[0] if not hist.loc[hist.index >= start_date_week.strftime('%Y-%m-%d')].empty else None
            price_month_ago = hist.loc[hist.index >= start_date_month.strftime('%Y-%m-%d')]['Close'].iloc[0] if not hist.loc[hist.index >= start_date_month.strftime('%Y-%m-%d')].empty else None
            
            change_week = ((current_price - price_week_ago) / price_week_ago) * 100 if price_week_ago else None
            change_month = ((current_price - price_month_ago) / price_month_ago) * 100 if price_month_ago else None
            
            market_cap_m = stock.info.get('marketCap', 0) / 1e6 if stock.info.get('marketCap') else None
            shares_outstanding = stock.info.get('sharesOutstanding', 0)
            
            data.append({
                "Company Name": company_name,
                "Ticker": symbol,
                "Current Price": current_price,
                "Price 1 Week Ago": price_week_ago,
                "Price 1 Month Ago": price_month_ago,
                "Change in Price (%) Week on Week": change_week,
                "Change in Price (%) 1 Month Ago": change_month,
                "Market Cap (M)": market_cap_m,
                "Number of Shares": shares_outstanding
            })
        else:
            print(f"No data found for {symbol}, it may be delisted or unavailable.")

    # Ensure data is returned
    if not data:
        print("No financial data was retrieved.")
    return data
          
def generate_ppt(slides_data, slides_data2):
    prs = Presentation()
    title_slide(prs)
    
    news_slide(prs, slides_data)
    news_slide(prs, slides_data2)
    
    financial_data = get_financial_data()
    comp_slide(prs, financial_data)
    
    commodities_slide(prs)
    project_slide(prs)
    grants_slide(prs)
    disclaimer_slide(prs)
    references_slide(prs)
    file_path = 'Monthly_Reports/Monthly_Report.pptx'
    prs.save(file_path)

def main():
    process = CrawlerProcess()
    process.crawl(Spider1)
    process.crawl(Spider2)
    process.start()
    
    slides_data = Spider1.slides_data
    slides_data2 = Spider2.slides_data
    
    generate_ppt(slides_data, slides_data2)
    
if __name__ == "__main__":
    main()