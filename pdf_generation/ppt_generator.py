from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Add a slide with a title layout
slide_layout = prs.slide_layouts[0]  # Use 0 for title slide layout
slide = prs.slides.add_slide(slide_layout)

# Set the title
title = slide.shapes.title
title.text = "Sample PowerPoint Presentation"

# Add bullet points
content = slide.placeholders[1]
content.text = "Bullet Points:"
content.text_frame.add_paragraph().text = "Point 1"
content.text_frame.add_paragraph().text = "Point 2"
content.text_frame.add_paragraph().text = "Point 3"

# Save the presentation
prs.save("sample_presentation.pptx")