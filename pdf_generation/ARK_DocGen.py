from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE 
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR 
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE
from datetime import datetime
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import os

def main():
    # Create a presentation object
    prs = Presentation()

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[5]  # Title slide layout
    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    title1.text = "ARK Energy"
    title1.top = Inches(2.5)
    title1.left = Inches(0.6)
    title1.width = Inches(2.66)
    title1.height = Inches(1.42)
    current_time= slide1.shapes.add_textbox(left=Inches(0.45), top=Inches(4.114), width=Inches(2.838), height=Inches(0.4055))
    current_datetime = datetime.now().strftime("%Y-%m-%d")
    # to generate time stamps use formating "%H:%M:%S" 
    current_time.text = f"Generated at: {current_datetime}"
    # Get the current directory of the Python script
    current_dir = os.path.dirname(os.path.abspath(__file__))
    img_folder = "Logo"
        
    # Path to the image file relative to the current directory

    img1 = "180degreelogo.png"
    img_path1 = os.path.join(current_dir, img_folder, img1)
    from_left = Inches(7.456)
    from_top = Inches(6.5)
    length = Inches(2)
    slide1.shapes.add_picture(img_path1, from_left, from_top, length)
    img2 = "Arkenergylogo.png"
    img_path2 = os.path.join(current_dir, img_folder, img2)
    from_left = Inches(5.169)
    from_top = Inches(6.5)
    length = Inches(2)
    slide1.shapes.add_picture(img_path2, from_left, from_top, length)

    # HOME BUTTON (OPTIONAL)
    home_button = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(8.311),top=Inches(0),width=Inches(0.708),height=Inches(0.606))
    fill_home = home_button.fill
    fill_home.solid()
    fill_home.fore_color.rgb = RGBColor(148, 197, 84)
    fill_home.zorder = 100
    line_home = home_button.line
    line_home.color.rgb = RGBColor(148, 197, 84)  
    line_home.zorder = 1


    img3 = "Home_button.png"
    img_path3 = os.path.join(current_dir, img_folder, img3)
    from_left = Inches(8.413)
    from_top = Inches(0.0511)
    height = Inches(0.5)
    width = Inches(0.5)
    add_picture3 = slide1.shapes.add_picture(img_path3, from_left, from_top, height, width)
    add_picture3.zorder = 100

    # Adding Divider Line
    green_line = slide1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.55), Inches(1), Inches(3.55), Inches(6.212))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)

    # # COMMODITY BUTTON
    # box_commodities = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    # fill_commodities = box_commodities.fill
    # fill_commodities.solid()
    # fill_commodities.fore_color.rgb = RGBColor(148, 197, 84)
    # fill_commodities.zorder = 1
    # line = box_commodities.line
    # line.color.rgb = RGBColor(148, 197, 84)  # HEX: 94C554
    # textbox_c = slide1.shapes.add_textbox(left=Inches(4.295), top=Inches(2.307), height=Inches(0.41), width=Inches(1.58))
    # text_box_c= textbox_c.text_frame
    # text_box_c.text = "Commodities"
    # box_commodities.click_action.target_slide = prs.slides[1] ## HERE

    # # NEWS BUTTON 
    # box_news = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    # fill_news = box_news.fill
    # fill_news.solid()
    # fill_news.fore_color.rgb = RGBColor(148, 197, 84)
    # fill_news.zorder = 1
    # line_news = box_news.line
    # line_news.color.rgb = RGBColor(148, 197, 84)  
    # textbox_n = slide1.shapes.add_textbox(left=Inches(6.464), top=Inches(2.311), height=Inches(0.41), width=Inches(0.8))
    # text_box_n= textbox_n.text_frame
    # text_box_n.text = "News"

    # # COMPETITORS/OMES BUTTON
    # box_comp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(7.877),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    # fill_comp = box_comp.fill
    # fill_comp.solid()
    # fill_comp.fore_color.rgb = RGBColor(148, 197, 84)
    # fill_comp.zorder = 2
    # line_comp = box_comp.line
    # line_comp.color.rgb = RGBColor(148, 197, 84) 
    # textbox_comp = slide1.shapes.add_textbox(left=Inches(7.944), top=Inches(2.314), height=Inches(0.41), width=Inches(1.527))
    # text_box_comp= textbox_comp.text_frame
    # text_box_comp.text = "Competitors"

    # # PROJECTS BUTTON
    # box_projects = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    # fill_projects = box_projects.fill
    # fill_projects.solid()
    # fill_projects.fore_color.rgb = RGBColor(148, 197, 84)
    # fill_projects.zorder = 2
    # line_projects = box_projects.line
    # line_projects.color.rgb = RGBColor(148, 197, 84) 
    # textbox_p = slide1.shapes.add_textbox(left=Inches(4.551), top=Inches(3.543), height=Inches(0.41), width=Inches(1.027))
    # text_box_p= textbox_p.text_frame
    # text_box_p.text = "Projects"

    # # GRANTS BUTTON
    # box_grants = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    # fill_grants = box_grants.fill
    # fill_grants.solid()
    # fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    # fill_grants.zorder = 2
    # line_grants = box_grants.line
    # line_grants.color.rgb = RGBColor(148, 197, 84) 
    # textbox_g = slide1.shapes.add_textbox(left=Inches(6.417), top=Inches(3.543), height=Inches(0.41), width=Inches(0.901))
    # text_box_g= textbox_g.text_frame
    # text_box_g.text = "Grants"

    # Slide 2: Content slide with bullet points
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "Key Commodities"
    title2.top = Inches(0.2)
    title2.left = Inches(0.3)
    title2.width = Inches(3)
    title2.height = Inches(1.42)
    title2.text_frame.paragraphs[0].font.size = Pt(26)

    # Divider
    green_line = slide2.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.55), Inches(0.5), Inches(3.55), Inches(7))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)

    # Main Body Table
    # Define the coordinates and size of the table
    x, y, cx, cy = Inches(0.4), Inches(1.3), Inches(2.8), Inches(1.8)
    # Add the table to the slide
    ammonia = slide2.shapes.add_table(9, 2, x, y, cx, cy)
    ammonia_table = ammonia.table
    # Set the fill color of the table to green
    for row in ammonia_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 

    # Title of the table
    cell = ammonia_table.cell(0, 0)
    cell.text = 'Ammonia'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Merge Relevant Cells
    merge_cell_ammonia = ammonia_table.cell(0, 0)
    other_cell_title = ammonia_table.cell(0, 1)
    merge_cell_ammonia.merge(other_cell_title)
    merge_cell_ammonia = ammonia_table.cell(8, 0)
    other_cell_title = ammonia_table.cell(8, 1)
    merge_cell_ammonia.merge(other_cell_title)

    # Body of the table
    current_datetime = "2024-04-30"
    labels = [
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = ammonia_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    # Merge last cell for the chart
    merge_cell = ammonia_table.cell(7, 0)
    other_cell = ammonia_table.cell(7, 1)
    merge_cell.merge(other_cell)

    # Last Row Larger 
    last_row = ammonia_table.rows[8]
    last_row.height = Inches(2)

    # SECOND TABLE 
    x, y, cx, cy = Inches(3.9), Inches(0.52), Inches(5.8), Inches(1.5)
    second = slide2.shapes.add_table(8, 3, x, y, cx, cy)
    second_table = second.table
    for row in second_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 
    cell = second_table.cell(0, 1)
    cell.text = 'Gold'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell = second_table.cell(0, 2)
    cell.text = 'Oil'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    current_datetime = "2024-04-30"
    labels = [
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = second_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # THIRD TABLE 
    x, y, cx, cy = Inches(3.9), Inches(3.8), Inches(5.8), Inches(1.5)
    third = slide2.shapes.add_table(8, 3, x, y, cx, cy)
    third_table = third.table
    for row in third_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 
    cell = third_table.cell(0, 1)
    cell.text = 'Steel'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell = third_table.cell(0, 2)
    cell.text = 'Zinc'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    current_datetime = "2024-04-30"
    labels = [
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = third_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Slide 3: Content slide with image
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide3 = prs.slides.add_slide(slide_layout)
    title3 = slide3.shapes.title
    title3.text = "News"

    # Defining the number of rows and columns for the table
    rows = 4
    cols = 3

    # Define the width of each column
    col_widths = [Inches(2)] * cols

    # Define the height of each row
    row_heights = [Inches(0.6), Inches(2.2), Inches(0.6), Inches(2.2)]

    # Adding a table 
    left = Inches(0.5)
    top = Inches(1.7)
    width = Inches(9)
    height = Inches(12)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # Set alternating fill colors for the rows
    for row_idx, row in enumerate(table.rows):
        if row_idx % 2 == 0:
            color = RGBColor(148, 197, 84)  # Green
            row.height = row_heights[0]
        else:
            color = RGBColor(255, 255, 255)  # White
            row.height = row_heights[1]
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = color

    # Add data to the table
    data = [
        ['A1', 'B1', 'C1'],
        ['A2', 'B2', 'C2'],
        ['A3', 'B3', 'C3'],
        ['A4', 'B4', 'C4']
    ]
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
            cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size
    
    # News Heading 
    cell = table.cell(0, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(0, 1)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(0, 2)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(2, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(2, 1)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(2, 2)
    cell.text = 'XX/XX/XXXX Title'
    
    # Content of News Articles 
    cell = table.cell(1, 0)
    cell.text = 'Content'
    cell = table.cell(1, 1)
    cell.text = 'Content'
    cell = table.cell(1, 2)
    cell.text = 'Content'
    cell = table.cell(3, 0)
    cell.text = 'Content'
    cell = table.cell(3, 1)
    cell.text = 'Content'
    cell = table.cell(3, 2)
    cell.text = 'Content'
    
    # Slide 4: Competitors Page
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide4 = prs.slides.add_slide(slide_layout)
    title4 = slide4.shapes.title
    title4.text = "Competitors/OEMs"
    title4.top = Inches(0.2)
    title4.left = Inches(0.3)
    title4.width = Inches(3)
    title4.height = Inches(1.42)
    title4.text_frame.paragraphs[0].font.size = Pt(25)
    
    # Divider
    green_line = slide4.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.55), Inches(0.5), Inches(3.55), Inches(7))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    
    # Main Body Table
    x, y, cx, cy = Inches(0.4), Inches(1.3), Inches(2.8), Inches(1.8)
    competitors = slide4.shapes.add_table(10, 2, x, y, cx, cy)
    competitors_table = competitors.table
    for row in competitors_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84)  # Green
    cell = competitors_table.cell(0, 0)
    cell.text = 'Company Name (Ticker: XXX)'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    merge_cell = competitors_table.cell(0, 0)
    other_cell = competitors_table.cell(0, 1)
    merge_cell.merge(other_cell)
    merge_cell = competitors_table.cell(8, 0)
    other_cell = competitors_table.cell(8, 1)
    merge_cell .merge(other_cell)
    merge_cell = competitors_table.cell(9, 0)
    other_cell = competitors_table.cell(9, 1)
    merge_cell .merge(other_cell)
    
    # Body of the table
    current_datetime = "2024-04-30"
    labels = [
        f"Primary Industry",
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = competitors_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(8)
    
    # Last Row Larger 
    last_row = competitors_table.rows[9]
    last_row.height = Inches(2)

    # SECOND TABLE 
    x, y, cx, cy = Inches(3.9), Inches(0.2), Inches(5.8), Inches(1.5)
    second = slide4.shapes.add_table(8, 3, x, y, cx, cy)
    second_table = second.table
    for row in second_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 
    cell = second_table.cell(0, 1)
    cell.text = 'Heavy Electrical Equipment'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    cell = second_table.cell(0, 2)
    cell.text = 'Automobile Manufactures'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    current_datetime = "2024-04-30"
    labels = [
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = second_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    # THIRD TABLE 
    x, y, cx, cy = Inches(3.9), Inches(3.6), Inches(5.8), Inches(1.5)
    third = slide4.shapes.add_table(8, 3, x, y, cx, cy)
    third_table = third.table
    for row in third_table.rows:
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(148, 197, 84) 
    cell = third_table.cell(0, 1)
    cell.text = 'Construction Machinery'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell = third_table.cell(0, 2)
    cell.text = 'Electric Components and Equipment'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    current_datetime = "2024-04-30"
    labels = [
        f"Current Price (as of {current_datetime})",
        f"Prev Wk Price (as of {current_datetime})",
        f"Price (as of {current_datetime})",
        "Change in Price (%)",
        f"Market Cap (as of {current_datetime})",
        "Number of Shares",
        "Chart (1 year)"
    ]
    for i, label in enumerate(labels):
        cell = third_table.cell(i + 1, 0)
        cell.text = label
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Slide 5: Projects
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide5 = prs.slides.add_slide(slide_layout)
    title5 = slide5.shapes.title
    title5.text = "Projects"
    title5.top = Inches(0.08)
    title5.left = Inches(0.45)
    title5.width = Inches(1.72)
    title5.height = Inches(1.42)
    title5.text_frame.paragraphs[0].font.size = Pt(28)

    rows = 6
    cols = 2
    
    # Define the width of each column
    col_widths = [Inches(2)] * cols
    
    # Define the height of each row
    row_heights = [Inches(1), Inches(1), Inches(1), Inches(1)]
   
    # Adding a table 
    left = Inches(0.45)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(12)
    table = slide5.shapes.add_table(rows, cols, left, top, width, height).table
   
    # Set alternating fill colors for the rows
    for row_idx, row in enumerate(table.rows):
        if row_idx % 2 == 0:
            color = RGBColor(148, 197, 84)  # Green
            row.height = row_heights[0]
        else:
            color = RGBColor(255, 255, 255)  # White
            row.height = row_heights[1]
        for cell in row.cells:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = color
    
    # Add data to the table
    data = [
        ['A1', 'B1'],
        ['A2', 'B2'],
        ['A3', 'B3'],
        ['A4', 'B4']
    ]
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
            cell.text_frame.paragraphs[0].font.size = Pt(12)  
    
    # Project Heading 
    cell = table.cell(0, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(1, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(2, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(3, 0)
    cell.text = 'XX/XX/XXXX Title'
    cell = table.cell(4, 0)
    cell.text = 'XX/XX/XXXX Title'
    
    # Content of Projects 
    cell = table.cell(0, 1)
    cell.text = 'Content'
    cell = table.cell(1, 1)
    cell.text = 'Content'
    cell = table.cell(2, 1)
    cell.text = 'Content'
    cell = table.cell(3, 1)
    cell.text = 'Content'
    cell = table.cell(4, 1)
    cell.text = 'Content'
    first = table.columns[0]
    first.width = Inches(2.5)
    second = table.columns[1]
    second.width = Inches(6.5)


    # Slide 6: Adding hyperlinks
    slide_layout = prs.slide_layouts[5] 
    slide6 = prs.slides.add_slide(slide_layout)
    title6 = slide6.shapes.title
    title6.text = "Grants"
    title6.top = Inches(0.08)
    title6.left = Inches(0.45)
    title6.width = Inches(1.72)
    title6.height = Inches(1.42)
    title6.text_frame.paragraphs[0].font.size = Pt(28)
   
    # Grant 1
    # Big White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(1.24),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)
    # Green Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(3.40),width=Inches(4.055),height=Inches(0.54))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(3.67),top=Inches(3.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text 
    textbox= slide6.shapes.add_textbox(left=Inches(0.84), top=Inches(1.4), height=Inches(0.41), width=Inches(1.58))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(3.67), top=Inches(3.468), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(1.22), Inches(0.7), Inches(3.94))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
    
    # GRANT 2 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.248),top=Inches(1.24),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.248),top=Inches(3.40),width=Inches(4.055),height=Inches(0.54))
    fill_apply = box_apply.fill
    fill_apply.solid()
    fill_apply.fore_color.rgb = RGBColor(148, 197, 84)
    fill_apply.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(8.2),top=Inches(3.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text
    textbox= slide6.shapes.add_textbox(left=Inches(5.38), top=Inches(1.37), height=Inches(0.41), width=Inches(1.58))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(8.2), top=Inches(3.46), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.24), Inches(1.22), Inches(5.24), Inches(3.94))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
   
    # GRANT 3 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(4.22),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box 
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(0.7),top=Inches(6.39),width=Inches(4.055),height=Inches(0.54))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(3.67),top=Inches(6.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Body Text
    textbox= slide6.shapes.add_textbox(left=Inches(0.84), top=Inches(4.53), height=Inches(1.31), width=Inches(3.91))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(3.67), top=Inches(6.448), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(4.22), Inches(0.7), Inches(6.93))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)
    
    # GRANT 4 
    # White Box
    box_grants = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.24),top=Inches(4.22),width=Inches(4.055),height=Inches(2.17))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 1
    line = box_grants.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Green Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(5.24),top=Inches(6.39),width=Inches(4.055),height=Inches(0.54))
    fill_apply = box_apply.fill
    fill_apply.solid()
    fill_apply.fore_color.rgb = RGBColor(148, 197, 84)
    fill_apply.zorder = 1
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Apply White Box
    box_apply = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,left=Inches(8.2),top=Inches(6.48),width=Inches(0.89),height=Inches(0.35))
    fill_grants = box_apply.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(255, 255, 255)
    fill_grants.zorder = 2
    line = box_apply.line
    line.color.rgb = RGBColor(0, 0, 0)  
    # Main Text Box
    textbox= slide6.shapes.add_textbox(left=Inches(5.38), top=Inches(4.53), height=Inches(1.31), width=Inches(3.91))
    text_box= textbox.text_frame
    text_box.text = "GRANT NAME"
    # Apply Button
    textbox= slide6.shapes.add_textbox(left=Inches(8.2), top=Inches(6.448), height=Inches(0.40), width=Inches(0.88))
    text_box= textbox.text_frame
    text_box.text = "APPLY"
    # Decorative Line
    green_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.24), Inches(4.2), Inches(5.24), Inches(6.93))
    line = green_line.line
    line.color.rgb = RGBColor(148, 197, 84)
    line.width = Inches(0.08)

    # Slide 7 Disclaimer
    slide_layout = prs.slide_layouts[5]  # Content slide layout
    slide7 = prs.slides.add_slide(slide_layout)
    title7 = slide7.shapes.title
    title7.text = "Disclaimer"
    numbered_list_items = [
        "Information Collection: The data collected through this project is sourced from publicly available websites and sources related to the ammonia and renewable energy sectors in Australia. It is intended to provide insights and analysis into these industries.",
        "Non-Commercial Use: The data and insights obtained from this project are not to be used for commercial purposes, including but not limited to selling or distributing the data for profit without explicit authorization.",
        "Accuracy and Reliability: While efforts are made to ensure the accuracy and reliability of the data collected, we do not guarantee its completeness or correctness. Users should independently verify the data before making any decisions or relying on it for business or investment purposes.",
        "No Endorsement: The inclusion of any specific data or information in this project does not constitute an endorsement or recommendation of any particular company, product, or service mentioned in the scraped data.",
        "No Liability: We disclaim any liability for damages, losses, or legal consequences arising from the use or misuse of the data collected through this project. Users assume all risks associated with their use of the information."
    ]
    # Add a text box for the numbered list
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6.5)
    text_box = slide7.shapes.add_textbox(left, top, width, height)
    # Set text alignment
    text_frame = text_box.text_frame
    text_frame.text_anchor = MSO_ANCHOR.TOP
    # Allow text to wrap within the text box
    text_frame.word_wrap = True
    # Add paragraphs with numbered text
    for idx, item in enumerate(numbered_list_items, start=1):
        p = text_frame.add_paragraph()
        p.text = f"{idx}. {item}"
        p.alignment = PP_ALIGN.LEFT
        # Set the paragraph to be numbered
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.level = 0  # Set the level of the numbered list (0 is the first level)
    # Add the disclaimer paragraph
    disclaimer = ("By using this program, you acknowledge that you have read, understood, "
                "and agree to be bound by this disclaimer.")
    p = text_frame.add_paragraph()
    p.text = disclaimer
    p.alignment = PP_ALIGN.LEFT

    # Slide 8 Adding hyperlinks
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide8 = prs.slides.add_slide(slide_layout)
    title8 = slide8.shapes.title
    title8.text = "References"
    para8 = slide8.placeholders[1].text_frame.paragraphs[0]
    addrun1 = para8.add_run()
    addrun1.text = "Google Hyperlink"
    hlink1 = addrun1.hyperlink
    hlink1.address = "https://www.google.com.au"
    # Save the presentation


        # COMMODITY BUTTON
    box_commodities = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_commodities = box_commodities.fill
    fill_commodities.solid()
    fill_commodities.fore_color.rgb = RGBColor(148, 197, 84)
    fill_commodities.zorder = 1
    line = box_commodities.line
    line.color.rgb = RGBColor(148, 197, 84)  # HEX: 94C554
    textbox_c = slide1.shapes.add_textbox(left=Inches(4.295), top=Inches(2.307), height=Inches(0.41), width=Inches(1.58))
    text_box_c= textbox_c.text_frame
    text_box_c.text = "Commodities"
    box_commodities.click_action.target_slide = prs.slides[1] ## HERE



    # NEWS BUTTON 
    box_news = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_news = box_news.fill
    fill_news.solid()
    fill_news.fore_color.rgb = RGBColor(148, 197, 84)
    fill_news.zorder = 1
    line_news = box_news.line
    line_news.color.rgb = RGBColor(148, 197, 84)  
    textbox_n = slide1.shapes.add_textbox(left=Inches(6.464), top=Inches(2.311), height=Inches(0.41), width=Inches(0.8))
    text_box_n= textbox_n.text_frame
    text_box_n.text = "News"
    box_news.click_action.target_slide = prs.slides[2]

    # COMPETITORS/OMES BUTTON
    box_comp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(7.877),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
    fill_comp = box_comp.fill
    fill_comp.solid()
    fill_comp.fore_color.rgb = RGBColor(148, 197, 84)
    fill_comp.zorder = 2
    line_comp = box_comp.line
    line_comp.color.rgb = RGBColor(148, 197, 84) 
    textbox_comp = slide1.shapes.add_textbox(left=Inches(7.944), top=Inches(2.314), height=Inches(0.41), width=Inches(1.527))
    text_box_comp= textbox_comp.text_frame
    text_box_comp.text = "Competitors"
    box_comp.click_action.target_slide = prs.slides[3]

    # PROJECTS BUTTON
    box_projects = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_projects = box_projects.fill
    fill_projects.solid()
    fill_projects.fore_color.rgb = RGBColor(148, 197, 84)
    fill_projects.zorder = 2
    line_projects = box_projects.line
    line_projects.color.rgb = RGBColor(148, 197, 84) 
    textbox_p = slide1.shapes.add_textbox(left=Inches(4.551), top=Inches(3.543), height=Inches(0.41), width=Inches(1.027))
    text_box_p= textbox_p.text_frame
    text_box_p.text = "Projects"
    box_projects.click_action.target_slide = prs.slides[4]

    # GRANTS BUTTON
    box_grants = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
    fill_grants = box_grants.fill
    fill_grants.solid()
    fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
    fill_grants.zorder = 2
    line_grants = box_grants.line
    line_grants.color.rgb = RGBColor(148, 197, 84) 
    textbox_g = slide1.shapes.add_textbox(left=Inches(6.417), top=Inches(3.543), height=Inches(0.41), width=Inches(0.901))
    text_box_g= textbox_g.text_frame
    text_box_g.text = "Grants"
    box_grants.click_action.target_slide = prs.slides[5]



    prs.save("ARK_Energy4.pptx")

    

if __name__ == "__main__":
    main()