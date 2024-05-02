from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE 
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR 
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Slide 1: Title slidex
slide_layout = prs.slide_layouts[5]  # Title slide layout
slide1 = prs.slides.add_slide(slide_layout)

title1 = slide1.shapes.title
title1.text = "ARK Energy"
title1.top = Inches(2.5)
title1.left = Inches(1)
title1.width = Inches(2.66)
title1.height = Inches(1.42)

current_time= slide1.shapes.add_textbox(left=Inches(0.9055), top=Inches(4.114), width=Inches(2.838), height=Inches(0.4055))
current_datetime = datetime.now().strftime("%Y-%m-%d")
# to generate time stamps use formating "%H:%M:%S" 
current_time.text = f"Generated at: {current_datetime}"

img1 = "180degreeslogo.png"
from_left = Inches(7.456)
from_top = Inches(6.5)
length = Inches(2)
add_picture = slide1.shapes.add_picture(img1, from_left, from_top, length)

img2 = "Arkenergylogo.png" 
from_left = Inches(5.169)
from_top = Inches(6.5)
length = Inches(2)
add_picture = slide1.shapes.add_picture(img2, from_left, from_top, length)

img3 = "Home_button.png" 
from_left = Inches(8.413)
from_top = Inches(0.0511)
height = Inches(0.5)
width = Inches(0.5)
add_picture3 = slide1.shapes.add_picture(img3, from_left, from_top, height, width)
add_picture3.zorder = 100

green_line = slide1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.921), Inches(1), Inches(3.921), Inches(6.212))
line = green_line.line
line.color.rgb = RGBColor(148, 197, 84)

# Home button 

home_button = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(8.311),top=Inches(0),width=Inches(0.708),height=Inches(0.606))
fill_home = home_button.fill
fill_home.solid()
fill_home.fore_color.rgb = RGBColor(148, 197, 84)
fill_home.zorder = 1
line_home = home_button.line
line_home.color.rgb = RGBColor(148, 197, 84)  
line_home.zorder = -16

# COMMODITY BUTTON

box_commodities = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
fill_commodities = box_commodities.fill
fill_commodities.solid()
fill_commodities.fore_color.rgb = RGBColor(148, 197, 84)
fill_commodities.zorder = 1
line = box_commodities.line
line.color.rgb = RGBColor(148, 197, 84)  # HEX: 94C554

textbox_c = slide1.shapes.add_textbox(left=Inches(4.295), top=Inches(2.307), height=Inches(0.41), width=Inches(1.58))
text_box_c= textbox_c.text_frame
text_box_c.text = "Commodities"

# NEWS BUTTON 

box_news = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
fill_news = box_news.fill
fill_news.solid()
fill_news.fore_color.rgb = RGBColor(148, 197, 84)
fill_news.zorder = 1
line_news = box_news.line
line_news.color.rgb = RGBColor(148, 197, 84)  

textbox_n = slide1.shapes.add_textbox(left=Inches(6.464), top=Inches(2.311), height=Inches(0.41), width=Inches(0.8))
text_box_n= textbox_n.text_frame
text_box_n.text = "News"

# COMPETITORS/OMES BUTTON

box_comp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(7.877),top=Inches(2.145),width=Inches(1.58),height=Inches(0.736))
fill_comp = box_comp.fill
fill_comp.solid()
fill_comp.fore_color.rgb = RGBColor(148, 197, 84)
fill_comp.zorder = 2
line_comp = box_comp.line
line_comp.color.rgb = RGBColor(148, 197, 84) 

textbox_comp = slide1.shapes.add_textbox(left=Inches(7.944), top=Inches(2.314), height=Inches(0.41), width=Inches(1.527))
text_box_comp= textbox_comp.text_frame
text_box_comp.text = "Competitors"

# PROJECTS BUTTON

box_projects = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(4.28), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
fill_projects = box_projects.fill
fill_projects.solid()
fill_projects.fore_color.rgb = RGBColor(148, 197, 84)
fill_projects.zorder = 2
line_projects = box_projects.line
line_projects.color.rgb = RGBColor(148, 197, 84) 

textbox_p = slide1.shapes.add_textbox(left=Inches(4.551), top=Inches(3.543), height=Inches(0.41), width=Inches(1.027))
text_box_p= textbox_p.text_frame
text_box_p.text = "Projects"

# GRANTS BUTTON

box_grants = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left=Inches(6.078), top=Inches(3.38), width=Inches(1.58), height=Inches(0.736))
fill_grants = box_grants.fill
fill_grants.solid()
fill_grants.fore_color.rgb = RGBColor(148, 197, 84)
fill_grants.zorder = 2
line_grants = box_grants.line
line_grants.color.rgb = RGBColor(148, 197, 84) 

textbox_g = slide1.shapes.add_textbox(left=Inches(6.417), top=Inches(3.543), height=Inches(0.41), width=Inches(0.901))
text_box_g= textbox_g.text_frame
text_box_g.text = "Grants"

# Slide 2: Content slide with bullet points
slide_layout = prs.slide_layouts[5]  # Content slide layout
slide2 = prs.slides.add_slide(slide_layout)
title2 = slide2.shapes.title
title2.text = "Key Commodities"
# Main Table 
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
ammonia = slide2.shapes.add_table(8, 2, x, y, cx, cy)
ammonia
ammonia.has_table 
ammonia_table = ammonia.table
ammonia_table
# Title of table 
cell = ammonia_table.cell(0, 0)
cell.text
cell.text = 'Ammonia'
# Merging cells 
merge_cell_ammonia = ammonia_table.cell(0, 0)
other_cell_title = ammonia_table.cell(0, 1)
merge_cell_ammonia.merge(other_cell_title)
merge_cell_ammonia.is_merge_origin
cell = ammonia_table.cell(1, 0)
# Body of Table 
cell.text
cell.text = f"Current Price (as of {current_datetime})"
cell = ammonia_table.cell(2, 0)
cell.text
cell.text = f"Prev Wk Price (as of {current_datetime})"
cell = ammonia_table.cell(3, 0)
cell.text
cell.text = f"Price (as of {current_datetime})"
cell = ammonia_table.cell(4, 0)
cell.text
cell.text = "Change in Price (%)"
cell = ammonia_table.cell(5, 0)
cell.text
cell.text = f"Market Cap (as of {current_datetime})"
cell = ammonia_table.cell(6, 0)
cell.text
cell.text = "Number of Shares"
cell = ammonia_table.cell(7, 0)
cell.text
cell.text = "Chart (1 year)"
merge_cell = ammonia_table.cell(7, 0)
other_cell = ammonia_table.cell(7, 1)
merge_cell.merge(other_cell)
merge_cell.is_merge_origin


# Slide 3: Content slide with image
slide_layout = prs.slide_layouts[5]  # Content slide layout
slide3 = prs.slides.add_slide(slide_layout)
title3 = slide3.shapes.title
title3.text = "News"

# Defining the number of rows and columns for the table
rows = 4
cols = 3

# Define the width of each column
col_widths = [Inches(2)] * cols

# Define the height of each row
row_heights = [Inches(0.8), Inches(1.5), Inches(0.8), Inches(1.5)]

# Adding a table 
left = Inches(0.5)
top = Inches(2)
width = Inches(9)
height = Inches(12)
table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

# Set alternating fill colors for the rows
for row_idx, row in enumerate(table.rows):
    if row_idx % 2 == 0:
        color = RGBColor(148, 197, 84)  # Green
        row.height = row_heights[0]
    else:
        color = RGBColor(255, 255, 255)  # White
        row.height = row_heights[1]
    for cell in row.cells:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = color

# Add data to the table
data = [
    ['A1', 'B1', 'C1'],
    ['A2', 'B2', 'C2'],
    ['A3', 'B3', 'C3'],
    ['A4', 'B4', 'C4']
]

for i, row in enumerate(data):
    for j, value in enumerate(row):
        cell = table.cell(i, j)
        cell.text = value
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
        cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size

# Slide 4: Content slide with image
slide_layout = prs.slide_layouts[5]  # Content slide layout
slide4 = prs.slides.add_slide(slide_layout)
title4 = slide4.shapes.title
title4.text = "Competitors/OEMs"

# Define the number of rows and columns for the table
rows = 8
cols = 2

# Define the width of each column
col_widths = [Inches(1)] * cols

# Define the height of each row
row_height = Inches(0.1)

# Add a table to the slide
left = Inches(0.6)
top = Inches(1.3)
width = Inches(4.5)
height = Inches(6)
competitor_table = slide4.shapes.add_table(rows, cols, left, top, width, height).table
cell.text = "Company Name (Ticker: XXX)"
cell = competitor_table.cell(0, 0)
cell.text = "Primary Industry"
cell = competitor_table.cell(1, 0)
cell.text = f"Current Price (as of {current_datetime})"
cell = competitor_table.cell(2, 0)
cell.text
cell.text = f"Prev Wk Price (as of {current_datetime})"
cell = competitor_table.cell(3, 0)
cell.text
cell.text = f"Price (as of {current_datetime})"
cell = competitor_table.cell(4, 0)
cell.text
cell.text = "Change in Price (%)"
cell = competitor_table.cell(5, 0)
cell.text
cell.text = f"Market Cap (as of {current_datetime})"
cell = competitor_table.cell(6, 0)
cell.text
cell.text = "Number of Shares"
cell = competitor_table.cell(7, 0)
cell.text
cell.text = "Chart (1 year)"  

#Merging Cells
merge_cell = competitor_table.cell(0, 0)
other_cell = competitor_table.cell(0, 1)
merge_cell.merge(other_cell)
merge_cell.is_merge_origin
cell = competitor_table.cell(1, 0)

merge_cell = competitor_table.cell(7, 0)
other_cell = competitor_table.cell(7, 1)
merge_cell.merge(other_cell)
merge_cell.is_merge_origin
cell = competitor_table.cell(7, 0)

# Set the fill color for the entire table
for row in competitor_table.rows:
    for cell in row.cells:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(148, 197, 84) 

# Add data to the table
for row in range(rows):
    for col in range(cols):
        cell = competitor_table.cell(row, col)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Center align text
        cell.text_frame.paragraphs[0].font.size = Pt(11)
           

# Slide 5: Adding hyperlinks
slide_layout = prs.slide_layouts[1]  # Content slide layout
slide5 = prs.slides.add_slide(slide_layout)
title5 = slide5.shapes.title
title5.text = "Hyperlinks"
para5 = slide5.placeholders[1].text_frame.paragraphs[0]
addrun1 = para5.add_run()
addrun1.text = "Google Hyperlink"
hlink1 = addrun1.hyperlink
hlink1.address = "https://www.google.com.au"

# addrun = object

# Save the presentation
prs.save("ARK_Energy.pptx")

