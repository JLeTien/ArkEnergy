from pptx import Presentation
from pptx.util import Inches
# Create a presentation object
prs = Presentation()

# Slide 1: Title slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Title Slide"

# Slide 2: Content slide with bullet points
slide_layout = prs.slide_layouts[1]  # Content slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Content Slide 1"
content = slide.placeholders[1]
content.text = "Bullet Points:"
content.text_frame.add_paragraph().text = "Point 1kehfblegfub"
content.text_frame.add_paragraph().text = "Point 2"
content.text_frame.add_paragraph().text = "Point 3"

# Slide 3: Content slide with image
slide_layout = prs.slide_layouts[1]  # Content slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Content Slide 2"
content = slide.placeholders[1]
content.text = "Image Slide"
#img_path = "path/to/your/image.jpg"
#content_pic = content.text_frame.add_paragraph().add_run()
#content_pic.add_picture(img_path, width=Inches(4), height=Inches(3))

# Save the presentation
prs.save("multiple_slides_presentation.pptx")