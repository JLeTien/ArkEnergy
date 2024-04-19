from pptx import Presentation
from pptx.util import Inches
# Create a presentation object
prs = Presentation()

# Slide 1: Title slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "ARK Energy"

# img1 = "180degreelogo.png"
from_left = Inches(7.3)
from_top = Inches(6.5)
length = Inches(2)
# add_picture = slide.shapes.add_picture(img1, from_left, from_top, length)

# img2 = "Arkenergylogo.png"
from_left = Inches(5)
from_top = Inches(6.5)
length = Inches(2)
# add_picture = slide.shapes.add_picture(img2, from_left, from_top, length)

# Slide 2: Content slide with bullet points
slide_layout = prs.slide_layouts[1]  # Content slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Content Slide 1"
content = slide.placeholders[1]
content.text = "Bullet Points:"
content.text_frame.add_paragraph().text = "Point 1"
content.text_frame.add_paragraph().text = "Point 2"
content.text_frame.add_paragraph().text = "Point 3"

# Slide 3: Content slide with image
slide_layout = prs.slide_layouts[1]  # Content slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Content Slide 2"


# Save the presentation
prs.save("multiple_slides_presentation.pptx")



